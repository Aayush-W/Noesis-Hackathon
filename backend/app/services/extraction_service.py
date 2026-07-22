import io
import re
import logging
import asyncio
from typing import Optional, List, Dict, Any

try:
    import cv2
except ImportError:
    cv2 = None

import httpx
import easyocr
import fitz  # PyMuPDF
import numpy as np
from PIL import Image
from app.core.config import settings

logger = logging.getLogger(__name__)

try:
    import docx
except ImportError:  # pragma: no cover - environment dependent
    docx = None

try:
    import pptx
except ImportError:  # pragma: no cover - environment dependent
    pptx = None

OCR_TEXT_THRESHOLD = 20
OCR_LANGUAGES = ["en"]
_ocr_reader: Optional[easyocr.Reader] = None
OCR_CONFIG_DEFAULT = {"detail": 1, "paragraph": True}
OCR_CONFIG_HANDWRITING = {
    "detail": 1,
    "paragraph": False,
    "contrast_ths": 0.05,
    "adjust_contrast": 0.7,
    "text_threshold": 0.4,
    "low_text": 0.2,
}
OCR_MIN_GOOD_TEXT = 220
OCR_MIN_GOOD_CONF = 0.50
OCR_MIN_GOOD_SCORE = 0.48


HF_TROCR_URL = "https://api-inference.huggingface.co/models/microsoft/trocr-base-handwritten"

def _get_ocr_reader() -> easyocr.Reader:
    """Lazily initialize EasyOCR reader to avoid startup overhead."""
    global _ocr_reader
    if _ocr_reader is None:
        _ocr_reader = easyocr.Reader(OCR_LANGUAGES, gpu=False)
    return _ocr_reader

def get_page_count(file_bytes: bytes) -> int:
    """Get total number of pages in a PDF document."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    count = len(doc)
    doc.close()
    return count

def rasterize_page(file_bytes: bytes, page_num: int, dpi: int = 300) -> np.ndarray:
    """Rasterize a PDF page to a numpy array (RGB)."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    page = doc.load_page(page_num)
    pix = page.get_pixmap(dpi=dpi)
    img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
    doc.close()
    return np.array(img)

def segment_lines(image_np: np.ndarray) -> List[np.ndarray]:
    """Segment an image into lines using horizontal projection profiling."""
    if cv2 is None:
        return [image_np]  # Fallback to full image if cv2 not available
        
    gray = cv2.cvtColor(image_np, cv2.COLOR_RGB2GRAY)
    # Thresholding to get binary image
    _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
    
    # Horizontal projection
    horizontal_projection = np.sum(binary, axis=1)
    
    # Find lines based on horizontal projection gaps
    height = horizontal_projection.shape[0]
    line_bounds = []
    in_line = False
    start = 0
    
    # Threshold for noise in projection
    proj_threshold = np.max(horizontal_projection) * 0.01
    
    for i in range(height):
        if not in_line and horizontal_projection[i] > proj_threshold:
            in_line = True
            start = i
        elif in_line and horizontal_projection[i] <= proj_threshold:
            in_line = False
            # Filter noise lines (e.g. underlines) shorter than 15px as per V2.3 plan
            if i - start >= 15:
                line_bounds.append((start, i))
                
    if in_line and height - start >= 15:
        line_bounds.append((start, height))
        
    # Crop lines
    line_crops = []
    for start, end in line_bounds:
        # Add a small padding
        p_start = max(0, start - 5)
        p_end = min(height, end + 5)
        line_crops.append(image_np[p_start:p_end, :])
        
    return line_crops

async def _call_hf_api(image_bytes: bytes) -> httpx.Response:
    """Call HuggingFace Inference API with authorization."""
    headers = {"Authorization": f"Bearer {settings.HF_TOKEN}"}
    async with httpx.AsyncClient(timeout=30) as client:
        return await client.post(HF_TROCR_URL, headers=headers, content=image_bytes)

async def run_trocr_inference(image_np: np.ndarray) -> str:
    """Run TrOCR inference on an image (page or line crop) with warmup handling."""
    if not settings.HF_TOKEN:
        logger.warning("HF_TOKEN missing; failing TrOCR path.")
        return ""
        
    # Convert NP array to bytes
    img = Image.fromarray(image_np)
    buf = io.BytesIO()
    img.save(buf, format="JPEG")
    img_bytes = buf.getvalue()
    
    # HF Cold Start / Warmup logic (3 retries)
    for attempt in range(3):
        try:
            response = await _call_hf_api(img_bytes)
            if response.status_code == 200:
                result = response.json()
                if isinstance(result, list) and len(result) > 0:
                    return result[0].get("generated_text", "")
                return ""
            elif response.status_code == 503:
                # Model loading
                wait_time = response.json().get("estimated_time", 20)
                logger.info(f"HuggingFace model loading. Waiting {wait_time}s (attempt {attempt+1}/3)")
                await asyncio.sleep(min(wait_time, 30))  # Cap wait time
                continue
            else:
                logger.error(f"HF API returned status {response.status_code}: {response.text}")
                return ""
        except Exception as e:
            logger.error(f"HF API connection failed: {e}")
            await asyncio.sleep(2)
            
    return ""

def is_low_confidence(text: str) -> bool:
    """Plausibility check for OCR results: >30% non-printable/non-ASCII chars = low confidence."""
    if not text or len(text) < 3:
        return True
    bad_chars = sum(1 for c in text if not c.isprintable() or ord(c) > 127)
    return (bad_chars / len(text)) > 0.30

async def extract_text_from_page(file_bytes: bytes, page_num: int) -> dict:
    """Async task to extract text from a single page using Adaptive TrOCR/Native strategy."""
    # 1. Start with native extraction (fast/digital)
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    page = doc.load_page(page_num)
    native_text = page.get_text().strip()
    doc.close()
    
    # 2. Heuristic: Sparse native text = likely scanned/handwritten
    if len(native_text) < OCR_TEXT_THRESHOLD or is_handwritten_page(native_text):
        logger.info(f"Starting Vision OCR for page {page_num + 1}...")
        
        # Adaptive DPI retry: 300 DPI first
        text = await _run_page_trocr(file_bytes, page_num, dpi=300)
        
        if is_low_confidence(text):
            logger.info(f"Low confidence text on page {page_num +1}, retrying at 600 DPI...")
            text = await _run_page_trocr(file_bytes, page_num, dpi=600)
            
        return {"text": text, "ocr_used": True}
        
    return {"text": native_text, "ocr_used": False}

async def _run_page_trocr(file_bytes: bytes, page_num: int, dpi: int) -> str:
    """Rasterize, segment, and run TrOCR for a full page."""
    image_np = rasterize_page(file_bytes, page_num, dpi=dpi)
    line_crops = segment_lines(image_np)
    
    line_texts = []
    for i, line_crop in enumerate(line_crops):
        try:
            # Wrap each line inference in timeout as per plan
            line_text = await asyncio.wait_for(
                run_trocr_inference(line_crop),
                timeout=60.0
            )
            if line_text:
                line_texts.append(line_text)
        except asyncio.TimeoutError:
            logger.warning(f"TrOCR timeout on page {page_num+1} line {i+1}")
            continue
            
    return "\n".join(line_texts)

def is_handwritten_page(text: str) -> bool:
    """Additional heuristic: high non-alpha ratio in sparse text suggest notes."""
    if not text: return True
    alpha_count = sum(1 for c in text if c.isalpha())
    return (alpha_count / len(text)) < 0.3


def _decode_text(file_bytes: bytes) -> str:
    """Decode text robustly with UTF-8 first, then latin-1 fallback."""
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1", errors="ignore")


def _run_easyocr(img_bytes: bytes) -> tuple[str, float]:
    """
    Run OCR on image bytes and return extracted text with average confidence.

    Returns:
        tuple[str, float]: (text, confidence)
    """
    reader = _get_ocr_reader()
    image = Image.open(io.BytesIO(img_bytes)).convert("RGB")
    image_np = np.array(image)
    gray_np = np.array(image.convert("L"))
    high_contrast_np = np.clip((gray_np.astype(np.float32) * 1.35), 0, 255).astype(np.uint8)

    def _normalize_text(text: str) -> str:
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    def parse_results(results: list) -> tuple[str, float]:
        if not results:
            return "", 0.0

        lines: list[str] = []
        confidences: list[float] = []

        for result in results:
            if len(result) < 3:
                continue
            text = str(result[1]).strip()
            conf = float(result[2])
            if text:
                lines.append(text)
                confidences.append(conf)

        if not lines:
            return "", 0.0

        avg_conf = sum(confidences) / len(confidences) if confidences else 0.0
        return _normalize_text("\n".join(lines)), avg_conf

    def _text_quality_score(text: str) -> float:
        if not text:
            return 0.0
        char_len = len(text)
        letters = sum(1 for ch in text if ch.isalpha())
        alpha_ratio = letters / max(char_len, 1)
        words = [w for w in re.split(r"\s+", text) if w]
        word_count = len(words)
        length_score = min(char_len / 800.0, 1.0)
        word_score = min(word_count / 120.0, 1.0)
        alpha_score = min(max(alpha_ratio, 0.0), 1.0)
        return 0.50 * length_score + 0.30 * alpha_score + 0.20 * word_score

    def _candidate_score(text: str, conf: float) -> float:
        quality = _text_quality_score(text)
        return 0.55 * conf + 0.45 * quality

    def _calibrate_confidence(text: str, conf: float) -> float:
        quality = _text_quality_score(text)
        calibrated = 0.45 * conf + 0.55 * quality
        return max(0.0, min(1.0, calibrated))

    def _run_candidate(img: np.ndarray, config: dict) -> tuple[str, float, float]:
        text, conf = parse_results(reader.readtext(img, **config))
        score = _candidate_score(text, conf)
        return text, conf, score

    # Handwriting-first fast path: high recall for scanned notes with lower latency.
    primary = _run_candidate(image_np, OCR_CONFIG_HANDWRITING)
    if (
        len(primary[0]) >= OCR_MIN_GOOD_TEXT
        and (primary[1] >= OCR_MIN_GOOD_CONF or primary[2] >= OCR_MIN_GOOD_SCORE)
    ):
        return primary[0], _calibrate_confidence(primary[0], primary[1])

    candidates: list[tuple[str, float, float]] = [primary]

    # Only attempt additional passes if primary result is weak.
    candidates.append(_run_candidate(gray_np, OCR_CONFIG_HANDWRITING))
    candidates.append(_run_candidate(high_contrast_np, OCR_CONFIG_HANDWRITING))
    candidates.append(_run_candidate(image_np, OCR_CONFIG_DEFAULT))

    # Heavy fallback variants for difficult handwritten pages.
    try:
        import cv2

        blur = cv2.GaussianBlur(gray_np, (3, 3), 0)
        adaptive = cv2.adaptiveThreshold(
            blur, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 31, 11
        )
        h, w = adaptive.shape
        upscaled = cv2.resize(adaptive, (int(w * 1.6), int(h * 1.6)), interpolation=cv2.INTER_CUBIC)

        candidates.append(_run_candidate(adaptive, OCR_CONFIG_HANDWRITING))
        candidates.append(_run_candidate(upscaled, OCR_CONFIG_HANDWRITING))
    except Exception:
        pass

    best = max(candidates, key=lambda x: x[2])
    if not best[0]:
        return "", 0.0

    return best[0], _calibrate_confidence(best[0], best[1])


async def extract_text_from_file(file_bytes: bytes, filename: str) -> dict:
    """Route file to the appropriate extraction method based on extension."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == "pdf":
        return await extract_pdf(file_bytes)
    if ext in ("png", "jpg", "jpeg"):
        return await extract_image(file_bytes)
    if ext == "docx":
        return extract_docx(file_bytes)
    if ext == "pptx":
        return extract_pptx(file_bytes)
    if ext == "txt":
        return {"text": _decode_text(file_bytes), "ocr_used": False, "confidence": 1.0}
    raise ValueError(f"Unsupported document format: {ext}")


async def extract_pdf(file_bytes: bytes) -> dict:
    """Extract text from PDF; use EasyOCR fallback for sparse/scanned pages."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages = []
    ocr_used = False
    ocr_confidences: list[float] = []

    ocr_pages_processed = 0
    MAX_OCR_PAGES = 20  # Limit to 20 pages of OCR to prevent 120s Timeout on CPUs

    for page_num, page in enumerate(doc):
        text = page.get_text().strip()

        # Sparse native text layer likely indicates scanned page.
        if len(text) < OCR_TEXT_THRESHOLD:
            if ocr_pages_processed < MAX_OCR_PAGES:
                ocr_used = True
                pix = page.get_pixmap(dpi=220)
                img_bytes = pix.tobytes("png")
                result = await extract_image(img_bytes, page_ref=f"p.{page_num + 1}")
                ocr_text = result.get("text", "").strip()
                ocr_confidences.append(float(result.get("confidence", 0.0)))

                # If OCR returns little/no text, preserve sparse native extraction.
                final_text = ocr_text if ocr_text else text
                pages.append({"page": page_num + 1, "text": final_text, "ocr": bool(ocr_text)})
                ocr_pages_processed += 1
            else:
                pages.append({"page": page_num + 1, "text": "[OCR Limit Reached - Page Skipped]", "ocr": False})
        else:
            pages.append({"page": page_num + 1, "text": text, "ocr": False})

    full_text = "\n\n".join(
        f"[Page {p['page']}]\n{p['text']}" for p in pages if p["text"].strip()
    )
    avg_conf = (sum(ocr_confidences) / len(ocr_confidences)) if ocr_confidences else 1.0

    return {
        "text": full_text,
        "ocr_used": ocr_used,
        "pages": pages,
        "confidence": avg_conf,
    }


async def extract_image(img_bytes: bytes, page_ref: str = "img") -> dict:
    """Extract text from images using EasyOCR."""
    import asyncio
    try:
        text, conf = await asyncio.to_thread(_run_easyocr, img_bytes)
        return {"text": text, "ocr_used": True, "confidence": conf, "pageRef": page_ref}
    except Exception as e:
        return {
            "text": "",
            "ocr_used": True,
            "confidence": 0.0,
            "pageRef": page_ref,
            "error": str(e),
        }


def extract_docx(file_bytes: bytes) -> dict:
    """Extract text from Word Document."""
    if docx is None:
        raise ValueError("python-docx is not installed; cannot parse DOCX files.")
    doc_io = io.BytesIO(file_bytes)
    doc = docx.Document(doc_io)
    text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    return {"text": text, "ocr_used": False, "confidence": 1.0}


def extract_pptx(file_bytes: bytes) -> dict:
    """Extract text from PowerPoint slides."""
    if pptx is None:
        raise ValueError("python-pptx is not installed; cannot parse PPTX files.")
    prs_io = io.BytesIO(file_bytes)
    prs = pptx.Presentation(prs_io)

    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_content.append(shape.text)
        if slide_content:
            slides_text.append(f"[Slide {i + 1}]\n" + "\n".join(slide_content))

    full_text = "\n\n".join(slides_text)
    return {"text": full_text, "ocr_used": False, "confidence": 1.0}
